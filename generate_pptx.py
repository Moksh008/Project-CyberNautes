import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette (Dark Cyber Theme)
    BG_DARK = RGBColor(10, 15, 26)       # #0a0f1a
    CARD_BG = RGBColor(18, 26, 43)       # #121a2b
    BLUE_ACCENT = RGBColor(59, 130, 246)  # #3b82f6
    PURPLE_ACCENT = RGBColor(168, 85, 247)# #a855f7
    CYAN_ACCENT = RGBColor(6, 182, 212)   # #06b6d4
    WHITE = RGBColor(255, 255, 255)
    GRAY_TEXT = RGBColor(156, 163, 175)   # #9ca3af
    RED_ACCENT = RGBColor(244, 63, 94)    # #f43f5e
    GREEN_ACCENT = RGBColor(52, 211, 153) # #34d399

    def apply_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.color.rgb = BG_DARK

    def add_header(slide, title_text, subtitle_text):
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(1.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = CYAN_ACCENT

    # -------------------------------------------------------------------------
    # SLIDE 1: TITLE SLIDE
    # -------------------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    apply_background(slide1)

    # Title Card Shape
    card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.2), Inches(10.933), Inches(5.1))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = BLUE_ACCENT
    card.line.width = Pt(1.5)

    tf1 = card.text_frame
    tf1.word_wrap = True
    tf1.margin_left = Inches(0.8)
    tf1.margin_top = Inches(0.8)
    tf1.margin_right = Inches(0.8)

    p = tf1.paragraphs[0]
    p.text = "SentinelAI"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT

    p_sub = tf1.add_paragraph()
    p_sub.text = "AI-Powered Cyber Defense Twin & Ephemeral Sandbox Platform"
    p_sub.font.size = Pt(22)
    p_sub.font.bold = True
    p_sub.font.color.rgb = WHITE
    p_sub.space_before = Pt(12)

    p_desc = tf1.add_paragraph()
    p_desc.text = "Autonomous threat reachability analysis, Red/Blue Team multi-agent reasoning, live container exploit detonation, and automated remediation code generation."
    p_desc.font.size = Pt(14)
    p_desc.font.color.rgb = GRAY_TEXT
    p_desc.space_before = Pt(18)

    p_footer = tf1.add_paragraph()
    p_footer.text = "Presented by Security Engineering Team • Capstone / MVP Technical Overview"
    p_footer.font.size = Pt(12)
    p_footer.font.color.rgb = CYAN_ACCENT
    p_footer.space_before = Pt(40)

    # -------------------------------------------------------------------------
    # SLIDE 2: PROBLEM STATEMENT & SOLUTION
    # -------------------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    apply_background(slide2)
    add_header(slide2, "1. Problem Statement & Core Value Proposition", "Moving security operations from passive vulnerability scanning to active twin-verified defense.")

    # Left Box: Problem
    box_p = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0))
    box_p.fill.solid()
    box_p.fill.fore_color.rgb = CARD_BG
    box_p.line.color.rgb = RED_ACCENT
    tf_p = box_p.text_frame
    tf_p.word_wrap = True
    tf_p.margin_left = tf_p.margin_right = tf_p.margin_top = Inches(0.4)

    p = tf_p.paragraphs[0]
    p.text = "The Problem in Modern SOCs"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RED_ACCENT

    bullets_p = [
        "Alert Fatigue: Scanners produce hundreds of unprioritized CVE alerts.",
        "Unknown Attack Reachability: Security teams don't know which CVEs create actual end-to-end attack paths to critical assets.",
        "Manual Remediation Bottlenecks: Analysts spend days deciding which mitigations are feasible without breaking production.",
        "Lack of Verification: Patches are deployed without proving they block real exploits beforehand."
    ]
    for b in bullets_p:
        p_b = tf_p.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = WHITE
        p_b.space_before = Pt(10)

    # Right Box: Solution
    box_s = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(5.0))
    box_s.fill.solid()
    box_s.fill.fore_color.rgb = CARD_BG
    box_s.line.color.rgb = GREEN_ACCENT
    tf_s = box_s.text_frame
    tf_s.word_wrap = True
    tf_s.margin_left = tf_s.margin_right = tf_s.margin_top = Inches(0.4)

    p = tf_s.paragraphs[0]
    p.text = "The SentinelAI Solution"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN_ACCENT

    bullets_s = [
        "Neo4j Knowledge Graph: Constructs a virtual Digital Twin connecting Assets ➔ Software ➔ CVEs ➔ MITRE Techniques.",
        "Shortest-Path Graph Reachability: Computes critical entry points (e.g. Internet ➔ Web ➔ App ➔ Database).",
        "Targeted Ephemeral Sandboxing: Detonates real red-team exploits in isolated Docker containers.",
        "AI Auto-Remediation: Generates deployable Bash scripts, Ansible playbooks, and Git diffs."
    ]
    for b in bullets_s:
        p_b = tf_s.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = WHITE
        p_b.space_before = Pt(10)

    # -------------------------------------------------------------------------
    # SLIDE 3: SYSTEM ARCHITECTURE & WORKFLOW
    # -------------------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    apply_background(slide3)
    add_header(slide3, "2. End-to-End System Architecture & Workflow", "A multi-layered pipeline from infrastructure ingestion to patch deployment.")

    steps = [
        ("Step 1: Ingestion", "Upload Infrastructure JSON or connect GitHub Repo to parse OS & services."),
        ("Step 2: Digital Twin", "Build Neo4j graph with relationships (Asset-RUNS-Software-HAS_CVE)."),
        ("Step 3: Graph Reachability", "Execute Cypher graph algorithms to find shortest attack paths to databases."),
        ("Step 4: Multi-Agent AI", "LangGraph agents: Offense AI analyzes threat, Defense AI suggests fixes."),
        ("Step 5: Ephemeral Sandbox", "Detonate Docker PoC exploits to prove before/after patch efficacy."),
        ("Step 6: Remediation Code", "Export production-ready Bash scripts, Ansible playbooks & Git diffs.")
    ]

    for idx, (title, desc) in enumerate(steps):
        col = idx % 3
        row = idx // 3
        x = Inches(0.8 + col * 4.0)
        y = Inches(1.8 + row * 2.6)
        
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(2.3))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = PURPLE_ACCENT if idx % 2 == 1 else BLUE_ACCENT
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.25)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = CYAN_ACCENT

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = WHITE
        p_d.space_before = Pt(8)

    # -------------------------------------------------------------------------
    # SLIDE 4: MVP SUPPORTED ATTACK SCENARIOS
    # -------------------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    apply_background(slide4)
    add_header(slide4, "3. Supported MVP Attack Scenarios", "Active red-team exploit detonations & verified mitigation in ephemeral containers.")

    scenarios = [
        ("Apache Web Server RCE", "CVE-2024-6387 / CVE-2021-41773", "Internet ➔ Apache Server ➔ RCE ➔ App ➔ DB", "Patch Apache, WAF Rule, Restrict DB Access"),
        ("ProFTPD Server Compromise", "CVE-2015-3306 (mod_copy execution)", "Internet ➔ ProFTPD FTP ➔ App ➔ DB", "Upgrade ProFTPD, Restrict FTP IPs, Disable Anon"),
        ("SQL Injection (SQLi) Web App", "OWASP Top 10 Injection", "Internet ➔ Web App ➔ SQLi Payload ➔ DB", "Prepared Statements, Input Sanitization, DB Least Privilege")
    ]

    for idx, (title, cve, chain, recs) in enumerate(scenarios):
        x = Inches(0.8 + idx * 4.0)
        y = Inches(1.8)
        
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(5.0))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BLUE_ACCENT
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)

        p = tf.paragraphs[0]
        p.text = f"Scenario #{idx+1}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = CYAN_ACCENT

        p_t = tf.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = WHITE
        p_t.space_before = Pt(4)

        p_cve = tf.add_paragraph()
        p_cve.text = "Vulnerability: " + cve
        p_cve.font.size = Pt(11)
        p_cve.font.color.rgb = RED_ACCENT
        p_cve.space_before = Pt(8)

        p_chain_head = tf.add_paragraph()
        p_chain_head.text = "Attack Chain:"
        p_chain_head.font.size = Pt(12)
        p_chain_head.font.bold = True
        p_chain_head.font.color.rgb = GRAY_TEXT
        p_chain_head.space_before = Pt(12)

        p_chain = tf.add_paragraph()
        p_chain.text = chain
        p_chain.font.size = Pt(11)
        p_chain.font.color.rgb = WHITE

        p_recs_head = tf.add_paragraph()
        p_recs_head.text = "AI Mitigations:"
        p_recs_head.font.size = Pt(12)
        p_recs_head.font.bold = True
        p_recs_head.font.color.rgb = GREEN_ACCENT
        p_recs_head.space_before = Pt(12)

        p_recs = tf.add_paragraph()
        p_recs.text = recs
        p_recs.font.size = Pt(11)
        p_recs.font.color.rgb = WHITE

    # -------------------------------------------------------------------------
    # SLIDE 5: TECH STACK, RESULTS & CONCLUSION
    # -------------------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    apply_background(slide5)
    add_header(slide5, "4. Technology Stack & Key Business Impact", "Full-stack implementation details and mathematical risk score reduction.")

    # Left: Tech Stack Table / Cards
    box_stack = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0))
    box_stack.fill.solid()
    box_stack.fill.fore_color.rgb = CARD_BG
    box_stack.line.color.rgb = PURPLE_ACCENT
    tf_st = box_stack.text_frame
    tf_st.word_wrap = True
    tf_st.margin_left = tf_st.margin_right = tf_st.margin_top = Inches(0.4)

    p = tf_st.paragraphs[0]
    p.text = "Production Technology Stack"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PURPLE_ACCENT

    stack_items = [
        ("Frontend", "React 19 + TypeScript + Tailwind CSS (Optimized)"),
        ("Backend API", "FastAPI (Python) + Pydantic v2"),
        ("Knowledge Graph", "Neo4j 5.x + Cypher Queries + NetworkX"),
        ("AI Agents", "LangGraph + LangChain + GPT-4o-mini / Gemini"),
        ("Threat Intel", "Live NIST NVD API v2 + MITRE ATT&CK Mapping"),
        ("Sandbox Engine", "Ephemeral Docker Engine API (Alpine Linux)")
    ]
    for cat, tech in stack_items:
        p_item = tf_st.add_paragraph()
        p_item.text = f"• {cat}: {tech}"
        p_item.font.size = Pt(12)
        p_item.font.color.rgb = WHITE
        p_item.space_before = Pt(8)

    # Right: Results & Impact
    box_imp = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(5.0))
    box_imp.fill.solid()
    box_imp.fill.fore_color.rgb = CARD_BG
    box_imp.line.color.rgb = CYAN_ACCENT
    tf_imp = box_imp.text_frame
    tf_imp.word_wrap = True
    tf_imp.margin_left = tf_imp.margin_right = tf_imp.margin_top = Inches(0.4)

    p = tf_imp.paragraphs[0]
    p.text = "Quantifiable Business Outcomes"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    impact_items = [
        ("Risk Score Reduction", "Proves mathematical drop in risk (e.g., 100 ➔ 0) upon verifying fixes."),
        ("Zero Production Risk", "All exploit detonations & patches run strictly in ephemeral containers."),
        ("Automated Code Export", "Instant generation of deployable Bash scripts, Ansible playbooks, and Git diffs."),
        ("Explainable AI Reports", "Produces Executive Summaries linking graph evidence to business impact.")
    ]
    for title, desc in impact_items:
        p_item = tf_imp.add_paragraph()
        p_item.text = f"✔ {title}: {desc}"
        p_item.font.size = Pt(12)
        p_item.font.color.rgb = WHITE
        p_item.space_before = Pt(10)

    # Save presentation
    output_filename = "SentinelAI_Presentation.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully to {output_filename}")

if __name__ == "__main__":
    create_presentation()
